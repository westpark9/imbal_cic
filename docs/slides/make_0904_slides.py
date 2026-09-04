#!/usr/bin/env python3
"""0904 랩미팅 발표자료 생성 스크립트 (KENTECH 템플릿 + docs/지침.txt) — 간략판(본문 슬라이드 9장, 표지·목차·섹션 없음, 사용자 편집본 0904_labmeeting_tabpfn2.pptx 피드백 반영).

    python docs/make_0904_slides.py            # docs/0904_labmeeting_tabpfn.pptx 생성 + 문체·기하 검사
    python docs/make_0904_slides.py --dump     # + 슬라이드별 텍스트 덤프

템플릿 `docs/KENTECH_발표자료_템플릿.pptx`의 TITLE_ONLY 레이아웃만 쓰고(우상단 로고 제거), 샘플 슬라이드 20장은 지운 뒤
샘플에서 읽어 둔 규격(카드 F3F6FA·헤더 E7F0F8·밴드 00306C 0.55in, 불릿 007BC6 75% 내어쓰기 0.2in, 표 헤더 네이비)으로 그린다.
수치는 랩로그 0902.md §8s~§8ag, run dir, budget_diag summary, main.tex, 서베이 xlsx에서 가져온 값이며 DATA 블록에 모아 둔다.
seed는 사용자 지시로 42~45 네 개만 표기한다(seed 46 붕괴·47·48은 슬라이드에서 제외, 랩로그 §8ad·§8ag 참조).
"""
import re, sys, os
from decimal import Decimal, ROUND_HALF_UP
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.dml import MSO_LINE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from lxml import etree

HERE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(HERE, 'KENTECH_발표자료_템플릿.pptx')
OUT = os.path.join(HERE, '0904_labmeeting_tabpfn.pptx')

NAVY, BLUE, CYAN = '00306C', '007BC6', '00B8EE'
CARD, CARD_HI, HEAD_HI = 'F3F6FA', 'E7F0F8', 'D6E4F3'
INK, GREY, GREY2, LINE = '21262E', '5D6773', '9AA3AE', 'D9E2EC'
OK, WARN = '1E7A46', 'B83A3A'
RAMP = ['00306C', '0A5490', '1173B2', '2585C4']

# ───────────────────────── DATA ─────────────────────────
def _r4(x): return float(Decimal(str(x)).quantize(Decimal('0.0001'), rounding=ROUND_HALF_UP))
SEEDS = ['42', '43', '44', '45']
_F1 = {  # global 열, 테스트 4,023,114행 (run 094650 / 122210 / 122328 / 174159)
    'benign':       [0.9909, 0.9924, 0.9878, 0.9927],
    'bot':          [0.9984, 0.9981, 0.9984, 0.9984],
    'brute_force':  [0.9159, 0.9159, 0.9158, 0.9158],
    'ddos':         [0.9459, 0.9588, 0.8879, 0.9596],
    'dos':          [0.7807, 0.7810, 0.7814, 0.7811],
    'infiltration': [0.4201, 0.4439, 0.4045, 0.4175],
    'web_attacks':  [0.2418, 0.2419, 0.2333, 0.2118],
    'macro':        [0.7562, 0.7617, 0.7442, 0.7539],
    'tail (dos·inf·web)': [0.5534, 0.5613, 0.5454, 0.5426],
}
_XGB = {  # (XGB 같은 100k C0 seed 42 = budget_diag_exp30_100k_075_bal, XGB 전체 풀 run 20260826_130310)
    'benign': (0.9924, 0.9943), 'bot': (0.9973, 0.9978), 'brute_force': (0.9159, 0.9159), 'ddos': (0.9994, 0.9921),
    'dos': (0.7818, 0.7822), 'infiltration': (0.3370, 0.4291), 'web_attacks': (0.2439, 0.1724),
    'macro': (0.7525, 0.7548), 'tail (dos·inf·web)': (0.5260, 0.5331),
}
def _avg(v): return _r4(sum(Decimal(str(x)) for x in v) / len(v))
PER_CLASS_F1 = {k: [f'{x:.4f}' for x in v] + [f'{_avg(v):.4f}', f'{_XGB[k][0]:.4f}', f'{_XGB[k][1]:.4f}'] for k, v in _F1.items()}
MEAN = _avg(_F1['macro']); TAIL = _avg(_F1['tail (dos·inf·web)'])
N_ABOVE = sum(1 for x in _F1['macro'] if x > 0.7548); N_TAIL_ABOVE = sum(1 for x in _F1['tail (dos·inf·web)'] if x > 0.5331)
GRID = [('100k', '0.600', '0.899', '0.7261', '0.7491'), ('250k', '0.836', '0.862', '0.6515 †', '0.7430'),
        ('500k', '0.580', '0.894', '0.7191', '0.7449'), ('1M', '0.656', '0.843', '0.7212', '0.7430')]
SCURVE = [('0.87 (자연)', '0.836', '0.9982', '0.705', '0.6515 †', 'dos'), ('0.75', '0.862', '0.9988', '0.813', '0.7430', 'brute'),
          ('0.60', '0.764', '0.9987', '0.092', '0.6366', 'inf'), ('0.50', '0.771', '0.9985', '0.776', '0.7343', 'brute'),
          ('0.40', '0.866', '0.9735', '0.131', '0.7180', 'brute')]
SAME_ROWS = [  # 같은 컨텍스트 행으로 XGBoost 재학습(scripts/results/budget_diag_exp30_*) vs TabPFN global(seed 42)
    ('자연 비율 · 250k', '217,676', '자연 (ddos 16,459 · inf 2,338)', '32', '0.7561', '0.6515 †'),
    ('benign 0.75 · 자연 배분 · 100k', '75,000', '자연 (ddos 12,730 · inf 1,809)', '24', '0.7572', '0.7491'),
    ('benign 0.75 · 균등 · 100k', '75,000', '4,848씩', '761', '0.7525', '0.7562'),
    ('benign 0.75 · 균등 · 250k', '187,500', '12,348씩', '761', '0.7474', '0.7513'),
    ('benign 0.75 · 자연 배분 · 1M', '750,000', '자연 (ddos 127,298 · inf 18,085)', '244', '0.7523', '0.7430'),
    ('전체 학습 풀 · 12.07M', '10,508,775', '전체', '1,522', '0.7548', '미실행'),
]

# ───────────────────────── 도형 헬퍼 ─────────────────────────
def rgb(h): return RGBColor.from_string(h)
def I(v): return Inches(v)

def rect(slide, x, y, w, h, fill=CARD, geom=MSO_SHAPE.ROUNDED_RECTANGLE, adj=None, line=None, line_w=None, dash=False):
    sh = slide.shapes.add_shape(geom, I(x), I(y), I(w), I(h)); sh.shadow.inherit = False
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill)
    if line is None: sh.line.fill.background()
    else:
        sh.line.color.rgb = rgb(line); sh.line.width = Pt(line_w or 1.25)
        if dash: sh.line.dash_style = MSO_LINE.DASH
    if adj is not None:
        for i, a in enumerate(adj if isinstance(adj, (list, tuple)) else [adj]): sh.adjustments[i] = a
    sh.text_frame.text = ''
    return sh

def _set_bullet(p, color=BLUE, size_pct=75000, char='•', marL=182880, indent=-182880):
    pPr = p._p.get_or_add_pPr(); pPr.set('marL', str(marL)); pPr.set('indent', str(indent))
    for tag in ('a:buClr', 'a:buSzPct', 'a:buFont', 'a:buChar', 'a:buNone'):
        for el in pPr.findall(qn(tag)): pPr.remove(el)
    buClr = etree.SubElement(pPr, qn('a:buClr')); etree.SubElement(buClr, qn('a:srgbClr')).set('val', color)
    etree.SubElement(pPr, qn('a:buSzPct')).set('val', str(size_pct))
    etree.SubElement(pPr, qn('a:buFont')).set('typeface', 'Arial')
    etree.SubElement(pPr, qn('a:buChar')).set('char', char)

def _no_bullet(p):
    pPr = p._p.get_or_add_pPr()
    for tag in ('a:buClr', 'a:buSzPct', 'a:buFont', 'a:buChar', 'a:buNone'):
        for el in pPr.findall(qn(tag)): pPr.remove(el)
    etree.SubElement(pPr, qn('a:buNone'))

def para(tf, runs, first=False, align='l', bullet=False, spc_before=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = {'l': PP_ALIGN.LEFT, 'c': PP_ALIGN.CENTER, 'r': PP_ALIGN.RIGHT}[align]
    p.space_before = Pt(spc_before)
    (_set_bullet if bullet else _no_bullet)(p)
    if isinstance(runs, str): runs = [(runs, 11.3, False, INK)]
    for t, sz, b, col in runs:
        r = p.add_run(); r.text = t; r.font.size = Pt(sz); r.font.bold = bool(b); r.font.color.rgb = rgb(col)
    return p

def _fill_tf(tf, paras, insets, anchor, wrap=True):
    tf.word_wrap = wrap; tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left, tf.margin_top, tf.margin_right, tf.margin_bottom = [Emu(v) for v in insets]
    tf.vertical_anchor = {'t': MSO_ANCHOR.TOP, 'm': MSO_ANCHOR.MIDDLE, 'b': MSO_ANCHOR.BOTTOM}[anchor]
    for i, pd in enumerate(paras):
        para(tf, pd['runs'], first=(i == 0), align=pd.get('align', 'l'), bullet=pd.get('bullet', False), spc_before=pd.get('spc', 0))

def textbox(slide, x, y, w, h, paras, insets=(0, 0, 0, 0), anchor='t', wrap=True):
    tb = slide.shapes.add_textbox(I(x), I(y), I(w), I(h)); _fill_tf(tb.text_frame, paras, insets, anchor, wrap); tb._insets = insets
    return tb

def shape_text(sh, paras, insets=(0, 0, 0, 0), anchor='m'):
    for pd in paras: pd.setdefault('align', 'c')
    _fill_tf(sh.text_frame, paras, insets, anchor); sh._insets = insets

def P(text, sz=11.3, bold=False, color=INK, **kw):
    d = dict(runs=[(text, sz, bold, color)]); d.update(kw); return d
def B(text, sz=11.3, color=INK, **kw):
    d = dict(runs=[(text, sz, False, color)], bullet=True); d.update(kw); return d
def R(runs, **kw):
    d = dict(runs=runs); d.update(kw); return d

CARD_INS = (182880, 80467, 146304, 65836)
HEAD_INS = (182880, 0, 146304, 0)

def card(slide, x, y, w, h, title, body_paras, head_fill=CARD_HI, body_fill=CARD, title_prefix=None, body_anchor='t'):
    rect(slide, x, y, w, h, fill=body_fill, adj=min(0.09 / max(h, 0.5), 0.05))
    hd = rect(slide, x, y, w, 0.44, fill=head_fill, geom=MSO_SHAPE.ROUND_2_SAME_RECTANGLE, adj=[0.18181, 0])
    runs = ([title_prefix] if title_prefix else []) + [(title, 12.5, True, NAVY)]
    shape_text(hd, [R(runs, align='l')], insets=HEAD_INS)
    return textbox(slide, x, y + 0.44, w, h - 0.44, body_paras, insets=CARD_INS, anchor=body_anchor)

def stat(slide, x, y, w, h, big, label, sub=None, fill=CARD, big_color=NAVY, big_sz=27):
    sh = rect(slide, x, y, w, h, fill=fill, adj=0.05161)
    ps = [P(big, big_sz, True, big_color), P(label, 10.5, True, INK)] + ([P(sub, 9.5, False, GREY)] if sub else [])
    shape_text(sh, ps, insets=(73152, 36576, 73152, 36576)); return sh

def arrow(slide, x, y, w=0.72, h=0.24, fill=BLUE, rot=0):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, I(x), I(y), I(w), I(h)); sh.shadow.inherit = False
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill); sh.line.fill.background(); sh.rotation = rot; return sh

def oval_num(slide, x, y, d, n, fill=NAVY, color='FFFFFF', sz=12.5):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, I(x), I(y), I(d), I(d)); sh.shadow.inherit = False
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill); sh.line.fill.background()
    shape_text(sh, [P(str(n), sz, True, color)], insets=(0, 0, 0, 0)); sh.text_frame.word_wrap = False; return sh

def chip(slide, x, y, w, h, text, fill=NAVY, color='FFFFFF', sz=9.5, bold=True):
    sh = rect(slide, x, y, w, h, fill=fill, adj=0.5)
    shape_text(sh, [P(text, sz, bold, color)], insets=(54864, 0, 54864, 0)); return sh

def band(slide, message, page):
    bd = rect(slide, 0, 6.95, 13.333, 0.55, fill=NAVY, geom=MSO_SHAPE.RECTANGLE)
    shape_text(bd, [P(message or '', 12.5, True, 'FFFFFF', align='l')], insets=(502920, 18288, 228600, 18288))
    textbox(slide, 12.42, 6.95, 0.60, 0.55, [P(str(page), 10, True, 'FFFFFF')], anchor='m', wrap=False)

def node(slide, x, y, w, h, title, sub=None, fill=CARD_HI, title_color=NAVY, sub_color=GREY, dash=False, line=None, tsz=11, ssz=9.5):
    sh = rect(slide, x, y, w, h, fill=fill, adj=0.12, line=line, line_w=1.25, dash=dash)
    ps = [P(title, tsz, True, title_color)] + ([P(sub, ssz, False, sub_color)] if sub else [])
    shape_text(sh, ps, insets=(54864, 27432, 54864, 27432)); return sh

# ───────────────────────── 텍스트 폭·줄수 추정 ─────────────────────────
def char_w(ch, sz):
    o = ord(ch)
    if ch == ' ': return 0.30 * sz
    if 0xAC00 <= o <= 0xD7A3 or 0x3000 <= o <= 0x9FFF or 0xFF00 <= o <= 0xFFEF: return 1.0 * sz
    if ch in '·•✓✕→↑↓×∪': return 0.9 * sz
    if ch.isdigit(): return 0.58 * sz
    if ch.isupper(): return 0.68 * sz
    if ch.isalpha(): return 0.54 * sz
    if ch in '.,:;\'"()[]|/': return 0.32 * sz
    return 0.55 * sz

def text_lines(s, width_in, sz):
    if not s: return 1
    avail = max(width_in * 72, 1); lines, cur = 1, 0.0
    for word in re.split(r'(\s+)', s):
        w = sum(char_w(c, sz) for c in word)
        if w > avail:
            for c in word:
                cw = char_w(c, sz)
                if cur + cw > avail: lines += 1; cur = cw
                else: cur += cw
            continue
        if cur + w > avail and word.strip(): lines += 1; cur = w
        else: cur += w
    return lines

def est_text_height(sh):
    ins = getattr(sh, '_insets', (0, 0, 0, 0)); w_in = Emu(sh.width).inches - (ins[0] + ins[2]) / 914400; h = 0.0
    for p in sh.text_frame.paragraphs:
        pPr = p._p.find(qn('a:pPr')); marL = int(pPr.get('marL', 0)) / 914400 if pPr is not None else 0
        txt = ''.join(r.text for r in p.runs); sz = max([r.font.size.pt for r in p.runs if r.font.size] or [11])
        n = text_lines(txt, w_in - marL, sz) if txt else 1
        h += n * sz * 1.25 / 72 + (p.space_before.pt if p.space_before else 0) / 72
    return h + (ins[1] + ins[3]) / 914400

def table(slide, x, y, colw, header, rows, hdr_sz=11, body_sz=10.5, row_h=0.36, hdr_h=0.40, hi_rows=(), bold_first_col=True,
          align=None, warn_cells=(), ok_cells=(), bold_cells=(), first_col_align='l'):
    nrows, ncols = len(rows) + 1, len(header)
    gf = slide.shapes.add_table(nrows, ncols, I(x), I(y), I(sum(colw)), I(hdr_h + row_h * len(rows))); tbl = gf.table
    tbl.first_row = True; tbl.horz_banding = False
    for i, w in enumerate(colw): tbl.columns[i].width = I(w)
    tbl.rows[0].height = I(hdr_h)
    align = align or (['l'] + ['c'] * (ncols - 1))
    def fill_cell(cell, text, sz, bold, color, bg, al):
        cell.fill.solid(); cell.fill.fore_color.rgb = rgb(bg)
        cell.margin_left, cell.margin_right, cell.margin_top, cell.margin_bottom = Emu(82296), Emu(64008), Emu(18288), Emu(18288)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE; tf = cell.text_frame; tf.word_wrap = True
        for li, line in enumerate(text.split('\n')):
            p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            p.alignment = {'l': PP_ALIGN.LEFT, 'c': PP_ALIGN.CENTER, 'r': PP_ALIGN.RIGHT}[al]; _no_bullet(p)
            rr = p.add_run(); rr.text = line; rr.font.size = Pt(sz); rr.font.bold = bold; rr.font.color.rgb = rgb(color)
    for c, htxt in enumerate(header): fill_cell(tbl.cell(0, c), htxt, hdr_sz, True, 'FFFFFF', NAVY, 'c' if c else first_col_align)
    est = hdr_h
    for r, row in enumerate(rows, start=1):
        bg = CARD_HI if (r - 1) in hi_rows else ('FFFFFF' if r % 2 == 1 else CARD); need = row_h
        for c, txt in enumerate(row):
            col = WARN if (r - 1, c) in warn_cells else (OK if (r - 1, c) in ok_cells else INK)
            bold = (c == 0 and bold_first_col) or ((r - 1, c) in bold_cells)
            fill_cell(tbl.cell(r, c), str(txt), body_sz, bold, col, bg, align[c] if c else first_col_align)
            lines = sum(text_lines(seg, colw[c] - (82296 + 64008) / 914400, body_sz) for seg in str(txt).split('\n'))
            need = max(need, lines * body_sz * 1.22 / 72 + 0.09)
        tbl.rows[r].height = I(need); est += need
    gf._est_h = est; return gf, est

# ───────────────────────── 프레젠테이션 준비 ─────────────────────────
prs = Presentation(TEMPLATE)
sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst): prs.part.drop_rel(sldId.rId); sldIdLst.remove(sldId)
LAYOUT = {l.name: l for l in prs.slide_layouts}
for sh in list(LAYOUT['TITLE_ONLY'].shapes):          # 우상단 학교 로고 제거 (사용자 지시)
    if sh.shape_type == 13: sh._element.getparent().remove(sh._element)

def set_title(slide, text):
    for ph in slide.placeholders:
        if ph.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE) or ph.placeholder_format.idx in (101, 102):
            ph.text_frame.text = text; return ph
    raise RuntimeError('title placeholder not found')

page = 0
def content_slide(title, subtitle=None):
    global page; page += 1
    s = prs.slides.add_slide(LAYOUT['TITLE_ONLY']); set_title(s, title)
    if subtitle: textbox(s, 0.63, 0.98, 12.09, 0.30, [P(subtitle, 11.5, False, GREY)])
    return s, page


# ═════════════ 1. Architecture ═════════════
s, pg = content_slide('Architecture', 'RACE-PFN · 오프라인 구성(위)과 온라인 추론(아래), 점선 상자는 학습형 게이트로 현재 보정 단계에서 닫혀 있음')
textbox(s, 0.62, 1.30, 6.0, 0.22, [P('OFFLINE · 한 번 구성한 뒤 고정', 9.5, True, GREY)])
off = [
    (0.62, 2.10, 'Global TFM  F(C0, ·)', 'C0 = benign 0.75 · 공격 균등 · 100k행'),
    (3.07, 2.20, '실패 서명', '위치 φ(16d) · 방향 y − p · 심각도, held-out 오류만'),
    (5.62, 2.00, '실패 영역 K개', '서명 공간 k-means, 서로소'),
    (7.97, 2.00, '전문가 은행 Cₖ', '균형 anchor ∪ 영역 행(≤186k)'),
    (10.32, 2.39, '상한(oracle)', '라벨로 행별 최선을 고르면 macro 0.83~0.87 > XGB'),
]
for i, (x, w, t, sub) in enumerate(off):
    node(s, x, 1.80, w, 1.0, t, sub, fill=CARD_HI if i < 4 else CARD, line=None if i < 4 else LINE)
for x, lab in [(2.72, 'held-out 오류'), (5.27, '군집'), (7.62, '컨텍스트화'), (9.97, '라벨로 평가')]:
    arrow(s, x + 0.04, 2.19, 0.27, 0.22, fill=GREY2)
    textbox(s, x - 0.45, 1.555, 1.25, 0.23, [P(lab, 8.5, False, GREY, align='c')])
textbox(s, 0.62, 3.30, 6.0, 0.26, [P('ONLINE · 입력마다', 9.5, True, GREY)])
on = [
    (0.62, 1.95, '1  Global 예측', '모든 입력, 한 번', False),
    (2.92, 2.30, '2  제안 게이트', 'P(이득 > 0) > τ_pre · 학습형 scorer', True),
    (5.57, 2.10, '3  top-1 전문가', '호출 1회, 컨텍스트 Cₖ로 예측', False),
    (8.02, 2.30, '4  검증 게이트', '이득 하한 > τ_post · 학습형 verifier', True),
    (10.67, 2.04, '교체 / 유지', '어느 게이트든 "아니오"면 global 유지', False),
]
for x, w, t, sub, dashed in on:
    node(s, x, 3.62, w, 1.0, t, sub, fill=CARD if dashed else CARD_HI, line=GREY2 if dashed else None, dash=dashed,
         title_color=GREY if dashed else NAVY)
for x in (2.57, 5.22, 7.67, 10.32):
    arrow(s, x + 0.04, 4.01, 0.27, 0.22, fill=GREY2)
textbox(s, 0.62, 4.95, 12.09, 0.6, [
    P('학습되는 것은 컨텍스트 행 집합(C0, Cₖ)과 게이트 모델 2개, backbone θ는 한 번도 갱신하지 않음 · 전문가는 공격 이름이 아니라 global의 실패 서명에서 자동으로 나옴', 10, False, GREY, align='c')])
band(s, None, pg)

# ═════════════ 2. TabPFN 소개 ═════════════
s, pg = content_slide('TabPFN · 학습 데이터를 입력으로 읽는 표 형식 foundation model')
textbox(s, 0.63, 0.98, 12.09, 0.30, [R([('Hollmann 외, ', 11.5, False, GREY), ('Accurate predictions on small data with a tabular foundation model', 11.5, True, GREY), (', Nature 2025 · 1세대는 ICLR 2023', 11.5, False, GREY)])])
rect(s, 0.62, 1.50, 6.55, 4.55, fill='FFFFFF', adj=0.02, line=LINE, line_w=1.25)
sh = rect(s, 0.92, 1.95, 1.95, 1.55, fill=CARD_HI, adj=0.08)
shape_text(sh, [P('컨텍스트', 11, True, NAVY), P('(x₁,y₁) … (xₙ,yₙ)', 10, False, INK), P('라벨 있는 행 n개', 9.5, False, GREY)], insets=(54864, 36576, 54864, 36576))
sh = rect(s, 0.92, 3.75, 1.95, 0.75, fill=CARD_HI, adj=0.12); shape_text(sh, [P('질의 행 x', 11, True, NAVY)], insets=(54864, 0, 54864, 0))
arrow(s, 2.97, 2.60, 0.55, 0.22); arrow(s, 2.97, 4.02, 0.55, 0.22)
sh = rect(s, 3.62, 2.35, 1.95, 2.15, fill=NAVY, adj=0.08)
shape_text(sh, [P('트랜스포머', 11.5, True, 'FFFFFF'), P('θ 고정', 11.5, True, CYAN), P('합성 데이터로 1회 사전학습', 9.5, False, 'D6E4F3')], insets=(54864, 36576, 54864, 36576))
arrow(s, 5.67, 3.31, 0.55, 0.22)
sh = rect(s, 6.32, 2.80, 0.72, 1.25, fill=CARD_HI, adj=0.08); shape_text(sh, [P('p(y|x)', 10.5, True, NAVY)], insets=(18288, 0, 18288, 0))
textbox(s, 0.92, 4.75, 5.95, 1.15, [
    P('경사 갱신 없음 · 적합 없음 · 순전파 한 번으로 테스트 전체 예측', 10, True, INK, align='c'),
    P('컨텍스트 행을 바꾸면 같은 θ로 다른 분류기가 됨, 전문화는 행 교체', 9.5, False, GREY, align='c')])
for i, (h1, h2) in enumerate([
    ('1  사전학습', '구조적 인과 모델(SCM) prior에서 뽑은 합성 데이터셋 약 1억 개로 트랜스포머를 한 번 학습, 실제 데이터는 쓰지 않음'),
    ('2  in-context 추론', '(학습 표본, 라벨) 시퀀스를 컨텍스트로 넣고 질의 행의 사후확률을 단일 순전파로 출력, 베이지안 추론 근사'),
    ('3  세대별 규모', '1세대 1K행 · 2세대 10K행×500특징 · 2.5세대 50K×2K · 3세대(2026.05) 100만 행, 이 연구는 3세대 분류 체크포인트 사용'),
    ('4  이 연구에서의 쓰임', '컨텍스트 100k~1M행을 실험 변수로, RTX 4090에서 4.02M행 추론에 100k 컨텍스트 기준 약 100초'),
]):
    y = 1.60 + 1.10 * i
    oval_num(s, 7.42, y + 0.10, 0.36, i + 1, sz=11.5)
    textbox(s, 7.95, y, 4.76, 1.0, [P(h1, 11.5, True, NAVY), P(h2, 10, False, GREY)])
band(s, None, pg)

# ═════════════ 3. 같은 행 비교 ═════════════
s, pg = content_slide('TabPFN 적합성 · 같은 컨텍스트 행으로 학습한 XGBoost와 비교')
hdr = ['컨텍스트 구성 (= XGBoost 학습 행)', 'benign', '공격 6종 (클래스당)', 'web', 'XGBoost', 'TabPFN-v3']
gf, est = table(s, 0.62, 1.45, [3.0, 1.35, 3.0, 0.8, 1.9, 2.04], hdr, [list(r) for r in SAME_ROWS], body_sz=10, hdr_sz=10.5,
                row_h=0.34, hdr_h=0.38, align=['l', 'r', 'l', 'r', 'c', 'c'], hi_rows=(2,), warn_cells={(0, 5)}, bold_cells={(2, 5), (5, 4)})
textbox(s, 0.62, 1.45 + est + 0.06, 12.09, 0.3, [P('† 특징이 완전히 같고 라벨만 다른 98k행이 dos로 예측되어 brute F1 0.49 · XGBoost 열은 각 run의 컨텍스트 행으로 같은 하이퍼파라미터로 재학습, TabPFN 열은 seed 42', 9.5, False, GREY)])
y2 = 1.45 + est + 0.42
card(s, 0.62, y2, 5.92, 6.05 - y2, 'XGBoost가 행 수·구성에 거의 무관한 이유 (0.747~0.757)', [
    B('남은 오류가 데이터 한계 · 특징이 완전히 같은데 라벨이 다른 98k행(ftp_bruteforce 77,344 + slowhttptest 21,110)은 어느 모델도 한쪽만 맞힘, dos F1 0.78이 상한', 10.5),
    B('infiltration은 benign과 분포가 겹쳐 recall 0.3대(XGBoost 0.32, TabPFN 0.27~0.32)', 10.5),
    B('그래서 100k만으로 이미 포화, 12.07M을 넣어도 0.7548', 10.5),
])
card(s, 6.79, y2, 5.92, 6.05 - y2, 'TabPFN이 구성에 따라 움직이는 이유 (0.65~0.76)', [
    B('in-context 학습은 컨텍스트의 클래스 비율을 prior로 읽음 · benign 87%면 tail이 눌려 0.65, benign 0.75 + 공격 균등이면 같은 행으로 XGBoost 이상(0.7562 vs 0.7525)', 10.5),
    B('곧 컨텍스트 구성이 재학습 없이 바꿀 수 있는 설계 변수, 구성 1개 비교 = 추론 1회(GPU 30분)', 10.5),
    B('대가 · 호출당 추론이 무거움(4.02M행 100초 GPU), 1M 컨텍스트는 19.5 GB이고 seed에 따라 불안정', 10.5),
], head_fill=CARD_HI)
band(s, None, pg)

# ═════════════ 4. SOTA 선정 ═════════════
s, pg = content_slide('SOTA 선정')
hdr = ['비교군', '계열', '분리하는 질문', '출처']
rows = [
    ('XGBoost (전체 풀 12.07M행)', '트리 부스팅', '표 형식 데이터의 강한 기본 기준선', 'KDD 2016'),
    ('TabPFN · 자연 비율 무작위 컨텍스트', '동일 모델 대조군', '구성 규칙(비중·배분) 자체의 효과를 분리', 'ICLR 2023 · Nature 2025'),
    ('BoostPFN', 'PFN 부스팅 앙상블', '잔차 기반 컨텍스트 앙상블이 단일 구성 컨텍스트를 넘는지', 'AISTATS 2025'),
    ('LoCalPFN · TabPFN-kNN', '질의별 검색 컨텍스트', '검색형(kNN) 컨텍스트가 규칙형 배분을 넘는지', 'NeurIPS 2024'),
    ('MixturePFN', '클러스터 라우팅 컨텍스트', '입력 공간 군집 라우팅이 실패 서명 군집(이 연구)과 어떻게 다른지', 'ICLR 2025'),
    ('DistPFN', '사후확률 라벨 시프트 보정', '컨텍스트를 바꾸지 않고 보정만으로 같은 이득이 나는지', 'ICML 2026'),
    ('Energy 기반 open-set flow 분류기', 'open-set NIDS', 'unseen 공격군 시나리오에서 energy 점수 기준선', 'Computers & Security 2025'),
]
gf, est = table(s, 0.62, 1.50, [3.2, 2.5, 4.4, 1.99], hdr, rows, body_sz=10.5, hdr_sz=11, row_h=0.46, hdr_h=0.40,
                align=['l', 'l', 'l', 'c'], hi_rows=(2,))
band(s, None, pg)

# ═════════════ 5. BoostPFN ═════════════
s, pg = content_slide('가장 가까운 논문 · BoostPFN')
textbox(s, 0.63, 0.98, 12.09, 0.30, [R([('Prior-Fitted Networks Scale to Larger Datasets When Treated as Weak Learners', 11.5, True, GREY), (', AISTATS 2025', 11.5, False, GREY)])])
rect(s, 0.62, 1.50, 6.55, 4.55, fill='FFFFFF', adj=0.02, line=LINE, line_w=1.25)
sh = rect(s, 0.92, 1.85, 1.75, 0.95, fill=CARD_HI, adj=0.1); shape_text(sh, [P('학습 데이터 D', 11, True, NAVY), P('표본 가중치 wₜ', 9.5, False, GREY)], insets=(36576, 18288, 36576, 18288))
arrow(s, 2.75, 2.21, 0.5, 0.22)
sh = rect(s, 3.33, 1.85, 1.95, 0.95, fill=CARD_HI, adj=0.1); shape_text(sh, [P('wₜ로 부분집합 추출', 10.5, True, NAVY), P('= 컨텍스트 Cₜ', 9.5, False, GREY)], insets=(36576, 18288, 36576, 18288))
arrow(s, 5.36, 2.21, 0.5, 0.22)
sh = rect(s, 5.94, 1.85, 1.0, 0.95, fill=NAVY, adj=0.1); shape_text(sh, [P('PFN fₜ', 10.5, True, 'FFFFFF'), P('θ 고정', 9, False, 'D6E4F3')], insets=(18288, 18288, 18288, 18288))
sh = rect(s, 3.33, 3.35, 3.61, 0.85, fill=CARD_HI, adj=0.1); shape_text(sh, [P('앙상블 Fₜ = Fₜ₋₁ + η·fₜ', 11, True, NAVY), P('잔차 rₜ = y − Fₜ(x) 계산', 9.5, False, GREY)], insets=(36576, 18288, 36576, 18288))
arrow(s, 6.33, 2.90, 0.22, 0.38, rot=90)
sh = rect(s, 0.92, 3.35, 2.05, 0.85, fill=CARD, adj=0.1, line=LINE, line_w=1); shape_text(sh, [P('가중치 갱신 wₜ₊₁', 10.5, True, NAVY), P('잔차 큰 표본 ↑', 9.5, False, GREY)], insets=(36576, 18288, 36576, 18288))
arrow(s, 3.05, 3.67, 0.22, 0.22, fill=GREY2, rot=180); arrow(s, 1.83, 2.90, 0.22, 0.38, fill=GREY2, rot=270)
textbox(s, 0.92, 4.45, 5.95, 1.45, [
    P('약한 학습기 하나 = 부분집합 컨텍스트를 가진 PFN 한 번의 추론,', 9.5, False, GREY, align='c'),
    P('잔차가 큰 표본이 다음 컨텍스트에 더 뽑히는 순차 앙상블', 9.5, False, GREY, align='c')])
for i, (h1, h2) in enumerate([
    ('1  PFN을 약한 학습기로', '학습 데이터 부분집합을 컨텍스트로 갖는 PFN 하나가 부스팅의 약한 학습기, 파라미터 갱신 없음'),
    ('2  잔차가 다음 컨텍스트를 정함', '현재 앙상블의 부스팅 잔차로 표본 가중치를 갱신하고 그 가중치로 새 컨텍스트를 추출해 순차 추가'),
    ('3  규모', '사전학습 컨텍스트 한계의 50배 데이터까지 표준 PFN을 상회, GBDT·딥러닝·AutoML보다 짧은 학습 시간(논문 보고)'),
    ('4  이 연구와의 차이', '잔차를 쓰는 점은 같으나 BoostPFN은 모든 입력에 모든 약한 학습기를 더하는 앙상블, 이 연구는 실패 영역별 컨텍스트를 이득이 기대되는 입력에만 top-1 호출'),
]):
    y = 1.60 + 1.10 * i
    oval_num(s, 7.42, y + 0.10, 0.36, i + 1, sz=11.5)
    textbox(s, 7.95, y, 4.76, 1.0, [P(h1, 11.5, True, NAVY), P(h2, 10, False, GREY)])
band(s, None, pg)

# ═════════════ 6. 실험 예정 ═════════════
s, pg = content_slide('실험 예정')
hdr = ['실험', '설정', '지표']
rows = [
    ('E1  closed-set 분류 · 4 데이터셋', 'benign 0.75 · 균등 · 100k 규칙을 ton_iot · bot_iot · unsw_nb15에 재조정 없이 적용, 시간순 분할, 노브는 검증 split(D_cal)에서 선택 후 test 1회, seed 전부 표기', 'macro-F1 · tail-macro-F1 · 클래스별 recall · benign recall'),
    ('E2  불균형 정도 sweep', '컨텍스트 풀의 tail 클래스 표본을 IR 10²·10³·10⁴로 줄여 재구성', 'tail-F1 대 IR 곡선과 심각도 AUC'),
    ('E3  unseen 공격군', '시간순 zero-shot novelty(시나리오 A) · unseen 수 sweep(시나리오 D), 논문 초안 5.2절', 'AUROC · FPR95 · AUOSCR(같은 known macro-F1에서) · benign 유지'),
    ('E4  비용', '컨텍스트 100k~1M, 추론 시간·GPU 메모리, 비교군도 같은 축', '초 · GB · macro-F1 대 컨텍스트 크기'),
]
gf, est = table(s, 0.62, 1.50, [2.6, 6.0, 3.49], hdr, rows, body_sz=10.5, hdr_sz=11, row_h=0.6, hdr_h=0.40, align=['l', 'l', 'l'])
band(s, None, pg)

# ═════════════ 7. Oracle과 실험결과 ═════════════
s, pg = content_slide('Oracle과 실험결과', '같은 split·같은 테스트 4,023,114행 macro-F1, 축은 0.60에서 시작 · oracle과 1M 자연 컨텍스트는 exp22c~24b(2026.08), 제안 모델은 exp31(2026.09)')
X0, X1, V0, V1 = 4.3, 12.5, 0.60, 0.90
def bx(v): return X0 + (v - V0) / (V1 - V0) * (X1 - X0)
bars = [
    ('전문가 은행 상한 (oracle)', '행마다 라벨로 global·전문가 중 최선을 고를 때', 0.83, 0.87, RAMP[0]),
    ('XGBoost 전체 풀', '12.07M행 학습', 0.7548, None, GREY2),
    ('RACE-PFN (제안) · 구성 컨텍스트 C0', 'benign 0.75 · 균등 · 100k, 4 seed 평균, 전문가 호출 전 단계', MEAN, None, RAMP[2]),
    ('TabPFN · 1M 자연 컨텍스트 (구성 전)', 'benign 87%, seed 42', 0.7212, None, RAMP[3]),
]
for i, (lab, sub, v, v2, col) in enumerate(bars):
    y = 1.70 + 0.62 * i
    textbox(s, 0.62, y - 0.04, 3.55, 0.55, [P(lab, 10.5, True, INK), P(sub, 9, False, GREY)])
    rect(s, X0, y + 0.08, X1 - X0, 0.28, fill='E4EAF1', geom=MSO_SHAPE.RECTANGLE)
    rect(s, X0, y + 0.08, bx(v) - X0, 0.28, fill=col, geom=MSO_SHAPE.RECTANGLE)
    if v2:
        rect(s, bx(v), y + 0.08, bx(v2) - bx(v), 0.28, fill='7F9CC4', geom=MSO_SHAPE.RECTANGLE)
        textbox(s, bx(v2) + 0.08, y + 0.05, 0.9, 0.32, [P(f'{v:.2f}~{v2:.2f}', 10.5, True, NAVY)], anchor='m')
    else:
        textbox(s, bx(v) + 0.08, y + 0.05, 0.9, 0.32, [P(f'{v:.4f}', 10.5, True, NAVY)], anchor='m')
for t in (0.60, 0.70, 0.80, 0.90):
    textbox(s, bx(t) - 0.3, 4.22, 0.6, 0.24, [P(f'{t:.2f}', 9, False, GREY2, align='c')])
card(s, 0.62, 4.75, 12.09, 1.30, 'oracle의 뜻', [
    B('테스트 행마다 라벨을 보고 global 예측과 전문가 K개의 예측 중 맞는 것을 고른 값, 전문가 은행이 도달할 수 있는 상한이며 실제 시스템의 성능이 아님', 10.5),
    B('상한이 기준선을 넘는다는 것은 전문가 컨텍스트에 정보가 있다는 뜻, 그 정보를 어느 입력에서 꺼내 쓸지가 다음 슬라이드들의 문제', 10.5),
])
band(s, None, pg)

# ═════════════ 8. 컨텍스트 구성 실험 ═════════════
s, pg = content_slide('컨텍스트 구성 실험 · 크기와 분포(자연 vs benign 비중) sweep', 'seed 42, 테스트 4,023,114행 macro-F1 · 빨간 수치의 원인은 아래 표')
hdr = ['크기', 'macro · 자연', 'macro · 0.75']
GRID_M = [(g[0], g[3].replace(' †', ''), g[4]) for g in GRID]
gf1, est1 = table(s, 0.83, 1.50, [1.15, 1.2, 1.2], hdr, GRID_M, body_sz=10.5, hdr_sz=10, row_h=0.36, hdr_h=0.40,
                  bold_cells={(0, 2), (1, 2), (2, 2), (3, 2)}, warn_cells={(1, 1)}, hi_rows=(0,))
textbox(s, 0.83, 1.50 + est1 + 0.08, 4.6, 0.5, [P('크기 sweep · benign 0.75는 100k~1M에서 0.743~0.749로 평평, 자연 분포는 250k에서 0.65로 꺼짐', 9.5, False, GREY)])
hdr = ['benign 비중 S', 'benign rec.', 'inf prec.', 'macro', '동일 벡터 98k행']
SC = [(r[0], r[2], r[3], r[4].replace(' †', ''), r[5]) for r in SCURVE]
gf2, est2 = table(s, 6.61, 1.50, [1.35, 1.05, 0.95, 0.85, 1.4], hdr, SC, body_sz=10.5, hdr_sz=10, row_h=0.36, hdr_h=0.40,
                  hi_rows=(1,), warn_cells={(0, 3), (0, 4), (2, 2), (2, 3), (2, 4), (4, 1), (4, 2)})
textbox(s, 6.61, 1.50 + est2 + 0.08, 5.6, 0.5, [P('분포 sweep(250k) · 0.75가 최적, 마지막 열은 동일 벡터 98k행이 어느 클래스로 예측되었나', 9.5, False, GREY)])
rect(s, 0.62, 4.30, 12.09, 2.55, fill=CARD_HI, adj=0.03)
hd = rect(s, 0.62, 4.30, 12.09, 0.44, fill=HEAD_HI, geom=MSO_SHAPE.ROUND_2_SAME_RECTANGLE, adj=[0.18181, 0])
shape_text(hd, [P('결과 해설 · 빨간 수치가 나온 이유 (동일 벡터 98k행 = ftp_bruteforce 77,344 + dos_slowhttptest 21,110행, 46개 특징이 같고 라벨만 다름)', 11.5, True, NAVY, align='l')], insets=HEAD_INS)
hdr = ['빨간 수치', '동일 벡터 98k행의 예측', '일어난 일', '결과']
rows = [
    ('자연 0.87 · macro 0.6515', 'dos', 'ftp_bruteforce 77,344행이 전부 dos로 (brute recall 0.33), dos는 FP 77k로 precision 0.44', 'brute F1 0.49 · dos F1 0.61, macro가 0.09 내려감'),
    ('S=0.60 · inf prec. 0.092', 'inf', '98k행이 inf로, inf FP 100,679행 중 98k가 이 블록', 'inf F1 0.14 · brute F1 0.49, macro 0.6366'),
    ('S=0.40 · benign rec. 0.9735 · inf prec. 0.131', 'brute (정상)', '원인이 다름 · benign 컨텍스트가 100k로 줄어 benign 90,251행이 inf로 누수', 'benign recall 0.025 하락, inf F1 0.19'),
    ('S=0.75 (채택)', 'brute (XGBoost와 같음)', '누수 없음, benign recall 0.9988', 'macro 0.7430, 이후 공격 균등 배분(web 61 → 761행)을 더한 것이 최종 구성'),
]
gf3, est3 = table(s, 0.82, 4.84, [2.6, 1.7, 4.3, 3.1], hdr, rows, body_sz=9.5, hdr_sz=10, row_h=0.36, hdr_h=0.34,
                  align=['l', 'c', 'l', 'l'], hi_rows=(3,), warn_cells={(0, 0), (1, 0), (2, 0)})
band(s, None, pg)

# ═════════════ 9. 클래스별 F1 ═════════════
s, pg = content_slide('100k 컨텍스트, benign 0.75, 공격 균등 · 클래스별 F1 (seed 42~45)')
hdr = ['클래스 F1'] + [f's{x}' for x in SEEDS] + ['4-seed 평균', 'XGB 같은 100k', 'XGB 전체 풀']
rows = [[k] + v for k, v in PER_CLASS_F1.items()]
gf, est = table(s, 0.62, 1.40, [2.2, 1.15, 1.15, 1.15, 1.15, 1.5, 1.85, 1.94], hdr, rows, body_sz=10, hdr_sz=10, row_h=0.30, hdr_h=0.36,
                hi_rows=(7,), warn_cells={(3, 3), (3, 5)}, bold_cells={(7, 5), (8, 5), (6, 5), (7, 7), (3, 6), (3, 7)})
y2 = 1.40 + est + 0.18
card(s, 0.62, y2, 5.92, 6.85 - y2, '클래스별 읽기', [
    B('macro 0.7540은 XGBoost 0.7548과 동률, tail(dos·inf·web) 0.551은 0.533을 4 seed 모두 상회', 10.5),
    B('DDoS 열세 · hoic 206k행 중 18k~53k행이 benign으로, 컨텍스트의 benign 추첨에 따라 seed별로 움직임(0.89~0.96)', 10.5),
    B('inf·web 우위 · 균등 배분으로 inf 4,847·web 761행이 컨텍스트에 들어가 XGBoost(inf 0.34, web 0.17)보다 높음', 10.5),
    B('dos·brute 동률 · 동일 벡터 98k행 때문에 어느 모델도 같은 상한, benign은 hoic 유출분만큼 소폭 열세', 10.5),
])
card(s, 6.79, y2, 5.92, 6.85 - y2, 'oracle 0.83~0.87과의 차이', [
    B('표의 값은 전문가 호출 전 단계 성능 · 제안 게이트가 보정 단계에서 닫혀 4 seed 중 3에서 호출 0건(seed 44만 열려 +0.002)', 10.5),
    B('게이트는 검증 창 예측으로 학습되는데 실패의 상당수가 test에만 나타남(SSH brute force 18:01 개시), 배울 표본이 없어 닫힘', 10.5),
    B('다음 실험 · 학습형 게이트를 무학습 규칙(p_benign < θ · D_cal 클래스 조건부 채택표)으로 바꿔 호출 활성화, 검증 split 기준 추첨 선택(exp34, 진행 중), 소량 라벨 컨텍스트 갱신', 10.5),
], head_fill=CARD_HI)
band(s, None, pg)

prs.save(OUT)
print('saved', OUT, 'slides', len(prs.slides))

# ───────────────────────── 검사 ─────────────────────────
def all_text_items(prs):
    for si, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    t = ''.join(r.text for r in p.runs)
                    if t.strip(): yield si, sh, p, t
            if sh.has_table:
                for r in sh.table.rows:
                    for c in r.cells:
                        for p in c.text_frame.paragraphs:
                            t = ''.join(rr.text for rr in p.runs)
                            if t.strip(): yield si, sh, p, t

problems = []
for si, sh, p, t in all_text_items(prs):
    if re.search(r'(합니다|입니다|습니다)(?![가-힣])', t): problems.append((si, '서술형 어미', t))
    if re.search(r'[—–]', t): problems.append((si, 'dash', t))
    if re.search(r'(?<![A-Za-z0-9])-|-(?![A-Za-z0-9])', t): problems.append((si, '하이픈', t))
    pPr = p._p.find(qn('a:pPr'))
    if pPr is not None and pPr.find(qn('a:buChar')) is not None and t.rstrip().endswith('.'): problems.append((si, '개조식 마침표', t))
    if re.search(r'^(출처|핵심)\s*:', t): problems.append((si, '라벨어', t))
    if '?' in t: problems.append((si, '질문형', t))

def box(sh): return (Emu(sh.left).inches, Emu(sh.top).inches, Emu(sh.left + sh.width).inches, Emu(sh.top + sh.height).inches)
def contains(a, b, tol=0.02): return a[0] <= b[0] + tol and a[1] <= b[1] + tol and a[2] >= b[2] - tol and a[3] >= b[3] - tol
def inter(a, b):
    w = min(a[2], b[2]) - max(a[0], b[0]); h = min(a[3], b[3]) - max(a[1], b[1]); return max(w, 0) * max(h, 0)
geo = []
for si, sl in enumerate(prs.slides, 1):
    shapes = list(sl.shapes)
    for sh in shapes:
        x0, y0, x1, y1 = box(sh)
        if x0 < -0.01 or y0 < -0.01 or x1 > 13.34 or y1 > 7.51: geo.append((si, '슬라이드 경계 밖', sh.name, round(x0, 2), round(y0, 2)))
        is_band_item = y0 >= 6.9 or (sh.shape_type == 1 and x0 == 0.0 and abs(x1 - 13.333) < 0.01)
        if y1 > 6.92 and not is_band_item: geo.append((si, '하단 밴드 침범', sh.name, round(y1, 2)))
        if sh.has_text_frame and sh.text_frame.text.strip() and sh.shape_type != 14:
            need = est_text_height(sh); h = y1 - y0
            if need > h * 0.99 and h > 0.2: geo.append((si, f'텍스트 넘침 추정 need={need:.2f} > box={h:.2f}', sh.text_frame.text[:50]))
    texty = [sh for sh in shapes if (sh.has_text_frame and sh.text_frame.text.strip() and sh.shape_type in (1, 17)) or sh.has_table]
    for a in range(len(texty)):
        for b in range(a + 1, len(texty)):
            A, Bx = box(texty[a]), box(texty[b])
            if inter(A, Bx) > 0.01 and not contains(A, Bx) and not contains(Bx, A):
                ta = texty[a].text_frame.text[:22] if texty[a].has_text_frame else 'TABLE'
                tb = texty[b].text_frame.text[:22] if texty[b].has_text_frame else 'TABLE'
                geo.append((si, '겹침', ta, tb))
print('\n=== 문체 검사:', len(problems), '건'); [print('  ', x) for x in problems]
print('=== 기하 검사:', len(geo), '건'); [print('  ', x) for x in geo]

if '--dump' in sys.argv:
    for si, sl in enumerate(prs.slides, 1):
        print(f'\n##### slide {si}'); seen = set()
        for sh in sl.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                t = sh.text_frame.text.replace('\n', ' / ')
                if t not in seen: seen.add(t); print('  ', t)
            if sh.has_table:
                for r in sh.table.rows: print('   |', ' | '.join(c.text.replace('\n', ' ') for c in r.cells))
