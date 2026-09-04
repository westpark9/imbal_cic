#!/usr/bin/env python3
"""0904_labmeeting_tabpfn5.pptx → 6.pptx: 서식·배치는 그대로, 한국어 문구만 영어로 치환.
같은 서식의 연속 run을 합친 뒤 문구 단위로 사전 치환한다. 사전에 없는 한글 문구는 경고로 출력."""
import re, sys, os
from pptx import Presentation
from pptx.util import Pt, Emu
HERE = os.path.dirname(os.path.abspath(__file__))
IN, OUT = os.path.join(HERE, '0904_labmeeting_tabpfn5.pptx'), os.path.join(HERE, '0904_labmeeting_tabpfn6.pptx')
T = {
# slide 1
'RACE-PFN · 오프라인 구성(위)과 온라인 추론(아래), 점선 상자는 학습형 게이트로 현재 보정 단계에서 닫혀 있음': 'RACE-PFN · offline construction (top) and online inference (bottom); dashed boxes are learned gates, currently closed at the calibration stage',
'OFFLINE · 한 번 구성한 뒤 고정': 'OFFLINE · built once, then frozen',
'C0 = benign 0.75 · 공격 균등 · 100k행': 'C0 = benign 0.75 · attacks equal · 100k rows',
'실패 서명': 'Failure signature',
'위치 φ(16d) · 방향 y − p · 심각도, held-out 오류만': 'position φ(16d) · direction y − p · severity, held-out errors only',
'실패 영역 K개': 'K failure regimes',
'서명 공간 k-means, 서로소': 'k-means in signature space, disjoint',
'전문가 은행 Cₖ': 'Expert bank Cₖ',
'균형 anchor ∪ 영역 행(≤186k)': 'balanced anchor ∪ regime rows (≤186k)',
'상한(oracle)': 'Upper bound (oracle)',
'라벨로 행별 최선을 고르면 macro 0.83~0.87 > XGB': 'label-picked best per row: macro 0.83~0.87 > XGB',
'held-out 오류': 'held-out errors', '군집': 'cluster', '컨텍스트화': 'to contexts', '라벨로 평가': 'score with labels',
'ONLINE · 입력마다': 'ONLINE · per input',
'1  Global 예측': '1  Global prediction', '모든 입력, 한 번': 'every input, once',
'2  제안 게이트': '2  Proposal gate', 'P(이득 > 0) > τ_pre · 학습형 scorer': 'P(gain > 0) > τ_pre · learned scorer',
'3  top-1 전문가': '3  Top-1 expert', '호출 1회, 컨텍스트 Cₖ로 예측': 'one call, predict with context Cₖ',
'4  검증 게이트': '4  Verification gate', '이득 하한 > τ_post · 학습형 verifier': 'gain lower bound > τ_post · learned verifier',
'교체 / 유지': 'Swap / keep', '어느 게이트든 "아니오"면 global 유지': '"no" at either gate keeps the global prediction',
'학습되는 것은 컨텍스트 행 집합(C0, Cₖ)과 게이트 모델 2개, backbone θ는 한 번도 갱신하지 않음 · 전문가는 공격 이름이 아니라 global의 실패 서명에서 자동으로 나옴':
    "Only the context row sets (C0, Cₖ) and the two gate models are learned; backbone θ is never updated · experts emerge from the global model's failure signatures, not from attack names",
# slide 2
'TabPFN · 학습 데이터를 입력으로 읽는 표 형식 foundation model': 'TabPFN · a tabular foundation model that reads training data as input',
'Hollmann 외,': 'Hollmann et al.,', ', Nature 2025 · 1세대는 ICLR 2023': ', Nature 2025 · 1st generation at ICLR 2023',
'컨텍스트': 'Context', '라벨 있는 행 n개': 'n labeled rows', '질의 행 x': 'query row x', '트랜스포머': 'Transformer', 'θ 고정': 'θ frozen',
'합성 데이터로 1회 사전학습': 'pre-trained once on synthetic data',
'경사 갱신 없음 · 적합 없음 · 순전파 한 번으로 테스트 전체 예측': 'No gradient step · no fitting · one forward pass predicts the whole test set',
'컨텍스트 행을 바꾸면 같은 θ로 다른 분류기가 됨, 전문화는 행 교체': 'Swapping context rows gives a different classifier with the same θ; specialization = row swap',
'1  사전학습': '1  Pre-training',
'구조적 인과 모델(SCM) prior에서 뽑은 합성 데이터셋 약 1억 개로 트랜스포머를 한 번 학습, 실제 데이터는 쓰지 않음': 'Transformer trained once on ~100M synthetic datasets drawn from a structural causal model (SCM) prior; no real data used',
'2  in-context 추론': '2  In-context inference',
'(학습 표본, 라벨) 시퀀스를 컨텍스트로 넣고 질의 행의 사후확률을 단일 순전파로 출력, 베이지안 추론 근사': "Feeds (sample, label) pairs as context and outputs the query row's posterior in one forward pass, approximating Bayesian inference",
'3  세대별 규모': '3  Scale by generation',
'1세대 1K행 · 2세대 10K행×500특징 · 2.5세대 50K×2K · 3세대(2026.05) 100만 행, 이 연구는 3세대 분류 체크포인트 사용': 'v1 1K rows · v2 10K rows × 500 features · v2.5 50K × 2K · v3 (2026.05) 1M rows; this work uses the v3 classification checkpoint',
'4  이 연구에서의 쓰임': '4  Role in this work',
'컨텍스트 100k~1M행을 실험 변수로, RTX 4090에서 4.02M행 추론에 100k 컨텍스트 기준 약 100초': 'Context of 100k~1M rows is the experimental variable; inference on 4.02M rows takes ~100 s on an RTX 4090 with a 100k context',
# slide 3
'TabPFN 적합성 · 같은 컨텍스트 행으로 학습한 XGBoost와 비교': 'Is TabPFN a fit · comparison with XGBoost trained on the same context rows',
'컨텍스트 구성 (= XGBoost 학습 행)': 'Context composition (= XGBoost training rows)', '공격 6종 (클래스당)': '6 attack classes (per class)',
'자연 비율 · 250k': 'Natural ratio · 250k', '자연 (ddos 16,459 · inf 2,338)': 'natural (ddos 16,459 · inf 2,338)',
'benign 0.75 · 자연 배분 · 100k': 'benign 0.75 · natural split · 100k', '자연 (ddos 12,730 · inf 1,809)': 'natural (ddos 12,730 · inf 1,809)',
'benign 0.75 · 균등 · 100k': 'benign 0.75 · equal · 100k', '4,848씩': '4,848 each',
'benign 0.75 · 균등 · 250k': 'benign 0.75 · equal · 250k', '12,348씩': '12,348 each',
'benign 0.75 · 자연 배분 · 1M': 'benign 0.75 · natural split · 1M', '자연 (ddos 127,298 · inf 18,085)': 'natural (ddos 127,298 · inf 18,085)',
'전체 학습 풀 · 12.07M': 'Full training pool · 12.07M', '전체': 'all',
'TabPFN이 구성에 따라 움직이는 이유 (0.65~0.76)': 'Why TabPFN moves with the composition (0.65~0.76)',
'in-context 학습은 컨텍스트의 클래스 비율을 prior로 읽음': 'In-context learning reads the class ratio of the context as a prior',
'대가 · 호출당 추론이 무거움(4.02M행 100초 GPU), 1M 컨텍스트는 19.5 GB이고 seed에 따라 불안정': 'Cost · heavy per-call inference (4.02M rows in 100 s on GPU); a 1M context needs 19.5 GB and is seed-unstable',
# slide 4
'SOTA 선정': 'SOTA selection', '비교군': 'Baseline', '계열': 'Family', '분리하는 질문': 'Question it isolates', '출처': 'Source',
'XGBoost (전체 풀 12.07M행)': 'XGBoost (full pool, 12.07M rows)', '트리 부스팅': 'Tree boosting', '표 형식 데이터의 강한 기본 기준선': 'Strong default baseline for tabular data',
'TabPFN · 자연 비율 무작위 컨텍스트': 'TabPFN · random natural-ratio context', '동일 모델 대조군': 'Same-model control',
'구성 규칙(비중·배분) 자체의 효과를 분리': 'Isolates the effect of the composition rule (share · allocation)',
'PFN 부스팅 앙상블': 'Boosted PFN ensemble', '잔차 기반 컨텍스트 앙상블이 단일 구성 컨텍스트를 넘는지': 'Does a residual-driven context ensemble beat one composed context',
'질의별 검색 컨텍스트': 'Per-query retrieved context', '검색형(kNN) 컨텍스트가 규칙형 배분을 넘는지': 'Does retrieval (kNN) context beat rule-based allocation',
'클러스터 라우팅 컨텍스트': 'Cluster-routed context', '입력 공간 군집 라우팅이 실패 서명 군집(이 연구)과 어떻게 다른지': 'How input-space clustering differs from failure-signature clustering (this work)',
'사후확률 라벨 시프트 보정': 'Posterior label-shift correction', '컨텍스트를 바꾸지 않고 보정만으로 같은 이득이 나는지': 'Does correction alone, without changing the context, give the same gain',
'Energy 기반 open-set flow 분류기': 'Energy-based open-set flow classifier', 'unseen 공격군 시나리오에서 energy 점수 기준선': 'Energy-score baseline for unseen-attack scenarios',
# slide 5
'학습 데이터 D': 'Training data D', '표본 가중치 wₜ': 'sample weights wₜ', 'wₜ로 부분집합 추출': 'Sample subset by wₜ', '= 컨텍스트 Cₜ': '= context Cₜ',
'앙상블 Fₜ = Fₜ₋₁ + η·fₜ': 'Ensemble Fₜ = Fₜ₋₁ + η·fₜ', '잔차 rₜ = y − Fₜ(x) 계산': 'residual rₜ = y − Fₜ(x)',
'가중치 갱신 wₜ₊₁': 'Update weights wₜ₊₁', '잔차 큰 표본 ↑': 'high-residual samples ↑',
'약한 학습기 하나 = 부분집합 컨텍스트를 가진 PFN 한 번의 추론,': 'One weak learner = one PFN inference with a subset context,',
'잔차가 큰 표본이 다음 컨텍스트에 더 뽑히는 순차 앙상블': 'a sequential ensemble where high-residual samples are drawn more into the next context',
'1  PFN을 약한 학습기로': '1  PFN as a weak learner',
'학습 데이터 부분집합을 컨텍스트로 갖는 PFN 하나가 부스팅의 약한 학습기, 파라미터 갱신 없음': 'A PFN whose context is a subset of the training data is the boosting weak learner; no parameter update',
'2  잔차가 다음 컨텍스트를 정함': '2  Residuals pick the next context',
'현재 앙상블의 부스팅 잔차로 표본 가중치를 갱신하고 그 가중치로 새 컨텍스트를 추출해 순차 추가': 'Boosting residuals of the current ensemble update the sample weights, which draw the next context, added sequentially',
'3  규모': '3  Scale',
'사전학습 컨텍스트 한계의 50배 데이터까지 표준 PFN을 상회, GBDT·딥러닝·AutoML보다 짧은 학습 시간(논문 보고)': 'Beats a plain PFN up to 50× the pre-training context limit, with shorter training time than GBDT, deep learning and AutoML (as reported)',
'4  이 연구와의 차이': '4  Difference from this work',
'잔차를 쓰는 점은 같으나 BoostPFN은 모든 입력에 모든 약한 학습기를 더하는 앙상블, 이 연구는 실패 영역별 컨텍스트를 이득이 기대되는 입력에만 top-1 호출': 'Both use residuals, but BoostPFN sums every weak learner on every input; this work calls one failure-regime context (top-1) only on inputs where a gain is expected',
# slide 6
'실험 예정': 'Planned experiments', '실험': 'Experiment', '설정': 'Setup', '지표': 'Metrics',
'E1  closed-set 분류 · 4 데이터셋': 'E1  Closed-set classification · 4 datasets',
'benign 0.75 · 균등 · 100k 규칙을 ton_iot · bot_iot · unsw_nb15에 재조정 없이 적용, 시간순 분할, 노브는 검증 split(D_cal)에서 선택 후 test 1회, seed 전부 표기': 'Apply the benign 0.75 · equal · 100k rule to ton_iot · bot_iot · unsw_nb15 without retuning, chronological split, knobs chosen on the validation split (D_cal) then one test run, all seeds reported',
'macro-F1 · tail-macro-F1 · 클래스별 recall · benign recall': 'macro-F1 · tail-macro-F1 · per-class recall · benign recall',
'E2  불균형 정도 sweep': 'E2  Imbalance-severity sweep',
'컨텍스트 풀의 tail 클래스 표본을 IR 10²·10³·10⁴로 줄여 재구성': 'Subsample tail classes in the context pool to IR 10²·10³·10⁴ and recompose',
'tail-F1 대 IR 곡선과 심각도 AUC': 'tail-F1 vs IR curve and severity AUC',
'E3  unseen 공격군': 'E3  Unseen attack families',
'시간순 zero-shot novelty(시나리오 A) · unseen 수 sweep(시나리오 D), 논문 초안 5.2절': 'Temporal zero-shot novelty (scenario A) · unseen-count sweep (scenario D), paper draft §5.2',
'AUROC · FPR95 · AUOSCR(같은 known macro-F1에서) · benign 유지': 'AUROC · FPR95 · AUOSCR (at matched known macro-F1) · benign retention',
'E4  비용': 'E4  Cost', '컨텍스트 100k~1M, 추론 시간·GPU 메모리, 비교군도 같은 축': 'Context 100k~1M, inference time · GPU memory, baselines on the same axes',
'초 · GB · macro-F1 대 컨텍스트 크기': 'seconds · GB · macro-F1 vs context size',
# slide 7
'Oracle과 실험결과': 'Oracle and results', '전문가 은행 상한 (oracle)': 'Expert-bank upper bound (oracle)',
'행마다 라벨로 global·전문가 중 최선을 고를 때': 'best of global/experts picked per row using labels',
'XGBoost 전체 풀': 'XGBoost, full pool', '12.07M행 학습': 'trained on 12.07M rows',
'제안 · 구성 컨텍스트 C0': 'Proposed · composed context C0',
'benign 0.75 · 균등 · 100k, 4 seed 평균, 전문가 호출 전 단계': 'benign 0.75 · equal · 100k, 4-seed mean, before expert calls',
'TabPFN · 1M 자연 컨텍스트 (구성 전)': 'TabPFN · 1M natural context (before composition)',
# slide 8
'100k 컨텍스트, benign 0.75, 공격 균등 · 클래스별 F1 (seed 42~45)': '100k context, benign 0.75, attacks equal · per-class F1 (seeds 42~45)',
'클래스 F1': 'Class F1', '4-seed 평균': '4-seed mean', 'XGB 같은 100k': 'XGB, same 100k', 'XGB 전체 풀': 'XGB, full pool', '결과': 'Results',
'macro 0.7540은 XGBoost 0.7548과 동률, tail(dos·inf·web) 0.551은 0.533을 4 seed 모두 상회': 'macro 0.7540 ties XGBoost 0.7548; tail (dos·inf·web) 0.551 > 0.533 in all 4 seeds',
'DDoS 열세나 아직 크게 개선되지 않은 inf, web에 대해서는 전문가 호출 게이트에 대한 추가 실험 필요.': 'DDoS deficit and still-modest inf/web gains: further experiments on the expert-call gate needed.',
}
hang = re.compile('[가-힣]')
def fmt(r):
    f = r.font; return (f.size.pt if f.size else None, f.bold, str(f.color.rgb) if (f.color is not None and f.color.type is not None) else None)
missing = set(); n_rep = 0
def translate_paragraph(p):
    global n_rep
    groups = []
    for r in p.runs:
        k = fmt(r)
        if groups and groups[-1][0] == k: groups[-1][1].append(r)
        else: groups.append([k, [r]])
    for _, runs in groups:
        text = ''.join(r.text for r in runs)
        if not hang.search(text): continue
        key = text.strip()
        if key in T:
            lead = text[:len(text) - len(text.lstrip())]; trail = text[len(text.rstrip()):]
            runs[0].text = lead + T[key] + trail; n_rep += 1
            for r in runs[1:]: r._r.getparent().remove(r._r)
        else:
            missing.add(text)
def all_tfs(prs):
    for sl in prs.slides:
        for sh in sl.shapes:
            if sh.has_text_frame: yield sh.text_frame
            if sh.has_table:
                for row in sh.table.rows:
                    for c in row.cells: yield c.text_frame
prs = Presentation(IN)
for tf in all_tfs(prs):
    for p in tf.paragraphs: translate_paragraph(p)
prs.save(OUT)
print('saved', OUT, 'replacements', n_rep)
if missing: print('MISSING:'); [print('   ', repr(m)) for m in sorted(missing)]

# ── 텍스트 여유 추정(영어는 한글보다 길어질 수 있음): 도형 실제 여백·글자 크기로 계산
def char_w(ch, sz):
    o = ord(ch)
    if ch == ' ': return 0.30 * sz
    if 0xAC00 <= o <= 0xD7A3 or 0x3000 <= o <= 0x9FFF: return 1.0 * sz
    if ch in '·•✓✕→↑↓×∪': return 0.9 * sz
    if ch.isdigit(): return 0.58 * sz
    if ch.isupper(): return 0.68 * sz
    if ch.isalpha(): return 0.54 * sz
    if ch in '.,:;\'"()[]|/': return 0.32 * sz
    return 0.55 * sz
def lines(s, width_in, sz):
    avail = max(width_in * 72, 1); n, cur = 1, 0.0
    for word in re.split(r'(\s+)', s):
        w = sum(char_w(c, sz) for c in word)
        if cur + w > avail and word.strip(): n += 1; cur = w
        else: cur += w
    return n
def need_h(sh):
    tf = sh.text_frame; ml, mr = (tf.margin_left or 0) / 914400, (tf.margin_right or 0) / 914400
    mt, mb = (tf.margin_top or 0) / 914400, (tf.margin_bottom or 0) / 914400
    w = Emu(sh.width).inches - ml - mr; h = 0.0
    for p in tf.paragraphs:
        pPr = p._p.find('{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        marL = int(pPr.get('marL', 0)) / 914400 if pPr is not None else 0
        txt = ''.join(r.text for r in p.runs); sz = max([r.font.size.pt for r in p.runs if r.font.size] or [11])
        h += (lines(txt, w - marL, sz) if txt else 1) * sz * 1.25 / 72 + (p.space_before.pt if p.space_before else 0) / 72
    return h + mt + mb
print('\n=== 텍스트 여유 추정 (need > box)')
for si, sl in enumerate(prs.slides, 1):
    for sh in sl.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() and not sh.is_placeholder and Emu(sh.height).inches > 0.2:
            need, box = need_h(sh), Emu(sh.height).inches
            if need > box * 0.99: print(f'  slide {si}: need={need:.2f} box={box:.2f}  {sh.text_frame.text[:70]!r}')
